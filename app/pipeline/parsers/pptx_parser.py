import structlog
from pptx import Presentation
from app.pipeline.parsers import ParsedChunk

logger = structlog.get_logger()

def parse_pptx(file_path: str) -> list[ParsedChunk]:
    """Parses a PPTX file and extracts text slide by slide."""
    logger.info("parser.pptx_started", path=file_path)
    chunks = []
    
    try:
        prs = Presentation(file_path)
        for i, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            
            if slide_text:
                combined_text = f"--- [SLIDE {i + 1}] ---\n\n" + "\n\n".join(slide_text)
                chunk = ParsedChunk(
                    text=combined_text,
                    page_number=i + 1,
                    chunk_index=i
                )
                chunks.append(chunk)
                
        logger.info("parser.pptx_finished", path=file_path, chunks=len(chunks))
        return chunks
    except Exception as e:
        logger.error("parser.pptx_failed", path=file_path, error=str(e))
        raise
