from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

# Create presentation object
prs = Presentation()

# Constants for styling
COLOR_PRIMARY = RGBColor(79, 70, 229)   # Indigo-600
COLOR_TEXT = RGBColor(30, 41, 59)      # Slate-800
COLOR_SECONDARY = RGBColor(100, 116, 139) # Slate-500
FONT_TITLE = 'Arial'
FONT_BODY = 'Calibri'

def add_title_slide(title, subtitle, name, company):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.color.rgb = COLOR_PRIMARY
    title_shape.text_frame.paragraphs[0].font.bold = True
    
    subtitle_shape = slide.placeholders[1]
    subtitle_shape.text = f"{subtitle}\n\n{name}\n{company}"
    subtitle_shape.text_frame.paragraphs[0].font.color.rgb = COLOR_TEXT

def add_content_slide(title, content_list):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.color.rgb = COLOR_PRIMARY
    
    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    tf.word_wrap = True
    
    for item in content_list:
        p = tf.add_paragraph()
        if isinstance(item, tuple):
            p.text = item[0]
            p.font.bold = True
            p.font.size = Pt(20)
            
            p2 = tf.add_paragraph()
            p2.text = item[1]
            p2.level = 1
            p2.font.size = Pt(16)
        else:
            p.text = item
            p.font.size = Pt(18)

def add_image_slide(title, img1_path, img2_path):
    slide_layout = prs.slide_layouts[1] # Content with title
    slide = prs.slides.add_slide(slide_layout)
    
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.color.rgb = COLOR_PRIMARY
    
    # Hide the default body placeholder
    slide.placeholders[1].visible = False
    
    # Add images side-by-side
    if os.path.exists(img1_path):
        slide.shapes.add_picture(img1_path, Inches(0.5), Inches(2), width=Inches(4.2))
    if os.path.exists(img2_path):
        slide.shapes.add_picture(img2_path, Inches(5.2), Inches(2), width=Inches(4.2))

# --- SLIDE 1: TITLE ---
add_title_slide(
    "DocuMind", 
    "AI-Powered Document Intelligence Platform using RAG Architecture",
    "Nilofar Sayyad (MCA II)",
    "Pirens IBMA, Loni"
)

# --- SLIDE 2: THE CHALLENGE ---
add_content_slide("The Challenge", [
    ("Information Overload", "Managing thousands of documents manually is impossible and prone to error."),
    ("Slow Retrieval", "Finding specific information in deeply nested PDFs takes hours of manual searching."),
    ("Context Blindness", "Traditional search fails to understand the actual meaning and context of queries.")
])

# --- SLIDE 3: THE SOLUTION ---
add_content_slide("The Solution: DocuMind", [
    "Leveraging Retrieval-Augmented Generation (RAG) to bridge the gap between static documents and AI.",
    ("Step 1: Ingestion", "Smart processing of PDFs and text files."),
    ("Step 2: Vectorization", "Encoding content into mathematical representations."),
    ("Step 3: Retrieval", "Context-aware search using pgvector."),
    ("Step 4: Generation", "Ultra-fast response generation via Groq Llama 3.")
])

# --- SLIDE 4: RAG 101 ---
add_content_slide("RAG 101: Core Concepts", [
    ("1. Chunks", "Breaking long documents into small, manageable pieces for the AI."),
    ("2. Embeddings", "Converting text into numerical vectors that represent concepts and meaning."),
    ("3. Retrieval", "Finding the exact chunks needed to answer a specific query."),
    ("4. Augmentation", "Giving the AI factual 'Notes' so it doesn't have to hallucinate.")
])

# --- SLIDE 5: PROJECT SCOPE ---
add_content_slide("Project Scope", [
    ("Target Audience", "Researchers, Academics, Legal Compliance Teams, and Corporate Knowledge Labs."),
    ("Semantic Intelligence", "Moving beyond keyword search to deep conceptual understanding of documents."),
    ("Scalability", "Architected to handle massive PDF libraries with near-instant retrieval speed.")
])

# --- SLIDE 6: INDUSTRY APPLICATIONS ---
add_content_slide("Real-World Industry Applications", [
    ("Legal Tech", "Automated contract auditing and fast retrieval of legal precedents."),
    ("Healthcare", "Scanning medical journals and patient records for fast diagnosis assistance."),
    ("Corporate HR", "Answering employee policy questions based on company handbooks."),
    ("Academic Research", "Analyzing thousands of research papers for literature reviews.")
])

# --- SLIDE 7: SYSTEM ARCHITECTURE ---
add_content_slide("System Architecture", [
    ("Frontend Layer", "React Client (Vite + Tailwind) for a responsive UI."),
    ("API Layer", "FastAPI Asynchronous Gateway for high-concurrency routing."),
    ("Logic Layer", "Celery & Redis for reliable background task processing."),
    ("Storage Layer", "PostgreSQL 16 + pgvector for high-speed similarity search."),
    ("AI Layer", "Groq LPU hardware for real-time Llama 3 inference.")
])

# --- SLIDE 8: TECHNICAL STACK ---
add_content_slide("Technical Stack", [
    "• Python 3.12: Core Backend Engine",
    "• PostgreSQL 16: Robust Database with pgvector extension",
    "• Docker: Containerized Service Orchestration",
    "• React.js: Premium Interactive Frontend",
    "• Groq API: Ultra-fast AI inference engine for source code analysis",
    "• HuggingFace: MiniLM-L6 embeddings",
    "• Redis: High-speed Message Broker",
    "• Prometheus: Real-time System Monitoring"
])

# --- SLIDE 10: PLATFORM CAPABILITIES ---
add_content_slide("Platform Capabilities", [
    ("Universal Code Intelligence", "Native support for C, C++, Java, Python, JS, & more."),
    ("Unified Viewer", "Preview PDFs, text, and code with live syntax highlighting."),
    ("Semantic Citations", "Every answer is backed by direct quotes and page references."),
    ("Vectorized Search", "Meaning-based retrieval beyond simple keywords."),
    ("Async OCR", "Handles batches of complex documents in the background."),
    ("Modern Security", "Built-in rate limiting and protected API endpoints.")
])

# --- SLIDE 11: USER MANAGEMENT MODULE ---
add_content_slide("User Management Module", [
    ("Secure Authentication", "JWT-based sessions with password hashing (BCrypt)."),
    ("Ownership Isolation", "Secure multi-tenancy where users only access their own data."),
    ("Member Profiles", "Personalized settings, theme management, and account dashboard."),
    ("Role-Based Flow", "Structured interactions for different user tiers (Member/Admin).")
])

# --- SLIDE 12: PROJECT WALKTHROUGH ---
add_image_slide("Project Walkthrough", "dashboard.png", "chat.png")

# --- SLIDE 13: FUTURE ROADMAP ---
add_content_slide("Future Roadmap", [
    ("Phase 1: Core RAG", "Database integration and async pipeline (Completed)."),
    ("Phase 2: Advanced Intelligence", "Multi-modal OCR support and vision-based interpretation."),
    ("Phase 3: Enterprise Auth", "Fine-grained ACL and role-based access control."),
    ("Phase 4: Collaboration", "Shared vector spaces and organization-level document pools.")
])

# --- SLIDE 14: ROADMAP DEEP-DIVE (NEW) ---
add_content_slide("Roadmap: Technical Implementation", [
    ("🖼️ Multi-modal OCR", "Using Tesseract/EasyOCR for scanning and Llama 3.2 Vision for diagrams."),
    ("🔐 Advanced Auth", "Expanding our JWT foundation to include OAuth2 and Multi-Factor Auth (MFA)."),
    ("👥 Team Collaboration", "Implementing Organisation-level Document Pools and Shared Vector Indexes."),
    ("🚀 Scalability", "Shifting to dedicated vector engines like Pinecone for billion-scale document handling.")
])

# --- SLIDE 15: CONCLUSION & IMPACT ---
add_content_slide("Conclusion & Impact", [
    ("Technical Achievement", "Distributed RAG pipeline with sub-second retrieval and generation."),
    ("Practical Utility", "Reduces document research time by up to 90%."),
    ("Platform Stability", "Scalable containerized architecture ready for production.")
])

# --- SLIDE 16: THANK YOU ---
add_title_slide("Thank You!", "Any Questions?", "DocuMind Project 2026", "Presented to Science & Tech Seminar")

# Save the presentation
output_path = "DocuMind_Presentation.pptx"
prs.save(output_path)
print(f"Presentation saved successfully as {output_path}")
